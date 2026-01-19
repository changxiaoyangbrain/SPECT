from pptx import Presentation
import sys
import os

def extract_text_from_pptx(file_path):
    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        return

    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error opening file: {e}")
        return

    with open("analysis_output.txt", "w", encoding="utf-8") as f:
        # Redirect stdout to file
        sys.stdout = f
        
        print(f"Analysis of Presentation: {os.path.basename(file_path)}\n")
        print(f"Total Slides: {len(prs.slides)}\n")

        for i, slide in enumerate(prs.slides):
            print(f"--- Slide {i+1} ---")
            # Try to find a title
            if slide.shapes.title:
                print(f"Title: {slide.shapes.title.text}")
            
            content = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # Skip title if we already printed it
                if shape == slide.shapes.title:
                    continue
                    
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        content.append(text)
            
            if content:
                print("Content:")
                for line in content:
                    print(f"- {line}")
            print("\n")
        
        # Reset stdout
        sys.stdout = sys.__stdout__
    print("Extraction complete. Output saved to analysis_output.txt")

if __name__ == "__main__":
    file_path = r"d:\SPECT\SPECT大作业.pptx"
    extract_text_from_pptx(file_path)
